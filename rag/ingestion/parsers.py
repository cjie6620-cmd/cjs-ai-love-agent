"""各文件格式的解析策略实现（标准库为主，PDF/Office 依赖可选包）。"""

from __future__ import annotations

import csv
import io
import json
import re
import zipfile
from html.parser import HTMLParser
from xml.etree.ElementTree import Element
from xml.etree.ElementTree import ParseError as XmlParseError
from xml.etree.ElementTree import fromstring as xml_fromstring

from .base import DocumentParser, ParsedDocument, ParseError


def _decode_best_effort(data: bytes) -> str:
    """目的：执行_decode_best_effort 方法相关逻辑。
    结果：返回当前步骤的处理结果，供后续流程继续使用。
    """
    for enc in ("utf-8", "utf-8-sig", "gbk", "latin-1"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("utf-8", errors="replace")


class PlainTextParser(DocumentParser):
    """目的：封装特定格式的解析逻辑，将原始内容转换为标准文本结果。
    结果：上游流程可以复用统一解析接口处理不同文件类型，降低适配成本。
    """

    # 目的：保存 SUFFIXES 字段，用于 PlainTextParser 的业务状态、配置或序列化。
    # 结果：实例在读写、校验和协作时可以获得稳定的 SUFFIXES 值。
    SUFFIXES = frozenset({".txt", ".text", ".log", ".md", ".markdown", ".rst"})

    def supports(self, mime: str | None, suffix: str) -> bool:
        """目的：根据当前输入执行条件判断，统一布尔分支的判定逻辑。
        结果：返回明确的判断结果，供上层决定后续流程。
        """
        if suffix in self.SUFFIXES:
            return True
        if mime and mime.startswith("text/") and mime not in ("text/html",):
            return True
        return False

    def parse(self, data: bytes, *, filename: str = "") -> ParsedDocument:
        """目的：将输入内容转换为统一的内部表示，屏蔽原始格式差异。
        结果：返回标准化解析结果，便于后续链路复用和扩展。
        """
        text = _decode_best_effort(data)
        return ParsedDocument(text=text, metadata={"parser": "plain", "filename": filename})


class HtmlParser(DocumentParser):
    """目的：封装特定格式的解析逻辑，将原始内容转换为标准文本结果。
    结果：上游流程可以复用统一解析接口处理不同文件类型，降低适配成本。
    """

    class _Stripper(HTMLParser):
        """目的：封装特定格式的解析逻辑，将原始内容转换为标准文本结果。
        结果：上游流程可以复用统一解析接口处理不同文件类型，降低适配成本。
        """
        def __init__(self) -> None:
            """目的：接收并保存运行所需的依赖、配置和初始状态。
            结果：实例初始化完成，可直接执行后续业务调用。
            """
            super().__init__()
            self.parts: list[str] = []

        def handle_data(self, data: str) -> None:
            """目的：封装当前步骤的核心处理逻辑，统一该能力的执行入口。
            结果：返回或落地稳定结果，供后续流程直接使用。
            """
            self.parts.append(data)

        def get_text(self) -> str:
            """目的：按指定条件读取目标数据、资源或结果集合。
            结果：返回可直接消费的查询结果，减少调用方重复处理。
            """
            return re.sub(r"\s+", " ", "".join(self.parts)).strip()

    # 目的：保存 SUFFIXES 字段，用于 HtmlParser 的业务状态、配置或序列化。
    # 结果：实例在读写、校验和协作时可以获得稳定的 SUFFIXES 值。
    SUFFIXES = frozenset({".html", ".htm", ".xhtml"})

    def supports(self, mime: str | None, suffix: str) -> bool:
        """目的：根据当前输入执行条件判断，统一布尔分支的判定逻辑。
        结果：返回明确的判断结果，供上层决定后续流程。
        """
        if suffix in self.SUFFIXES:
            return True
        return bool(mime and mime.startswith("text/html"))

    def parse(self, data: bytes, *, filename: str = "") -> ParsedDocument:
        """目的：将输入内容转换为统一的内部表示，屏蔽原始格式差异。
        结果：返回标准化解析结果，便于后续链路复用和扩展。
        """
        raw = _decode_best_effort(data)
        stripper = HtmlParser._Stripper()
        stripper.feed(raw)
        text = stripper.get_text()
        return ParsedDocument(text=text, metadata={"parser": "html", "filename": filename})


class JsonParser(DocumentParser):
    """目的：封装特定格式的解析逻辑，将原始内容转换为标准文本结果。
    结果：上游流程可以复用统一解析接口处理不同文件类型，降低适配成本。
    """

    # 目的：保存 SUFFIXES 字段，用于 JsonParser 的业务状态、配置或序列化。
    # 结果：实例在读写、校验和协作时可以获得稳定的 SUFFIXES 值。
    SUFFIXES = frozenset({".json", ".jsonl"})

    def supports(self, mime: str | None, suffix: str) -> bool:
        """目的：根据当前输入执行条件判断，统一布尔分支的判定逻辑。
        结果：返回明确的判断结果，供上层决定后续流程。
        """
        if suffix in self.SUFFIXES:
            return True
        return mime in ("application/json", "application/x-ndjson")

    def parse(self, data: bytes, *, filename: str = "") -> ParsedDocument:
        """目的：将输入内容转换为统一的内部表示，屏蔽原始格式差异。
        结果：返回标准化解析结果，便于后续链路复用和扩展。
        """
        raw = _decode_best_effort(data)
        try:
            obj = json.loads(raw)
        except json.JSONDecodeError as e:
            raise ParseError("invalid json") from e
        if isinstance(obj, (dict, list)):
            text = json.dumps(obj, ensure_ascii=False, indent=2)
        else:
            text = str(obj)
        return ParsedDocument(text=text, metadata={"parser": "json", "filename": filename})


class CsvParser(DocumentParser):
    """目的：封装特定格式的解析逻辑，将原始内容转换为标准文本结果。
    结果：上游流程可以复用统一解析接口处理不同文件类型，降低适配成本。
    """

    # 目的：保存 SUFFIXES 字段，用于 CsvParser 的业务状态、配置或序列化。
    # 结果：实例在读写、校验和协作时可以获得稳定的 SUFFIXES 值。
    SUFFIXES = frozenset({".csv", ".tsv"})
    # 目的：保存 DELIM 字段，用于 CsvParser 的业务状态、配置或序列化。
    # 结果：实例在读写、校验和协作时可以获得稳定的 DELIM 值。
    DELIM = {".csv": ",", ".tsv": "\t"}

    def supports(self, mime: str | None, suffix: str) -> bool:
        """目的：根据当前输入执行条件判断，统一布尔分支的判定逻辑。
        结果：返回明确的判断结果，供上层决定后续流程。
        """
        if suffix in self.SUFFIXES:
            return True
        return mime in ("text/csv", "text/tab-separated-values")

    def parse(self, data: bytes, *, filename: str = "") -> ParsedDocument:
        """目的：将输入内容转换为统一的内部表示，屏蔽原始格式差异。
        结果：返回标准化解析结果，便于后续链路复用和扩展。
        """
        raw = _decode_best_effort(data)
        ext = filename[filename.rfind(".") :].lower() if "." in filename else ".csv"
        delim = self.DELIM.get(ext, ",")
        lines: list[str] = []
        reader = csv.reader(io.StringIO(raw), delimiter=delim)
        for row in reader:
            lines.append(" | ".join(cell.strip() for cell in row))
        text = "\n".join(lines)
        return ParsedDocument(text=text, metadata={"parser": "csv", "filename": filename})


class XmlParser(DocumentParser):
    """目的：封装特定格式的解析逻辑，将原始内容转换为标准文本结果。
    结果：上游流程可以复用统一解析接口处理不同文件类型，降低适配成本。
    """

    # 目的：保存 SUFFIXES 字段，用于 XmlParser 的业务状态、配置或序列化。
    # 结果：实例在读写、校验和协作时可以获得稳定的 SUFFIXES 值。
    SUFFIXES = frozenset({".xml", ".svg"})

    def supports(self, mime: str | None, suffix: str) -> bool:
        """目的：根据当前输入执行条件判断，统一布尔分支的判定逻辑。
        结果：返回明确的判断结果，供上层决定后续流程。
        """
        if suffix in self.SUFFIXES:
            return True
        return bool(mime and ("xml" in mime))

    def parse(self, data: bytes, *, filename: str = "") -> ParsedDocument:
        """目的：将输入内容转换为统一的内部表示，屏蔽原始格式差异。
        结果：返回标准化解析结果，便于后续链路复用和扩展。
        """
        try:
            root = xml_fromstring(data)
        except XmlParseError as e:
            raise ParseError("invalid xml") from e

        def collect(el: Element) -> list[str]:
            """目的：把当前节点、子节点和尾随文本按遍历顺序抽取出来。
            结果：返回去空白后的文本片段列表。
            """
            out: list[str] = []
            if el.text and el.text.strip():
                out.append(el.text.strip())
            for child in el:
                out.extend(collect(child))
            if el.tail and el.tail.strip():
                out.append(el.tail.strip())
            return out

        text = "\n".join(collect(root))
        return ParsedDocument(text=text, metadata={"parser": "xml", "filename": filename})


class PdfParser(DocumentParser):
    """目的：封装特定格式的解析逻辑，将原始内容转换为标准文本结果。
    结果：上游流程可以复用统一解析接口处理不同文件类型，降低适配成本。
    """

    # 目的：保存 SUFFIXES 字段，用于 PdfParser 的业务状态、配置或序列化。
    # 结果：实例在读写、校验和协作时可以获得稳定的 SUFFIXES 值。
    SUFFIXES = frozenset({".pdf"})

    def supports(self, mime: str | None, suffix: str) -> bool:
        """目的：根据当前输入执行条件判断，统一布尔分支的判定逻辑。
        结果：返回明确的判断结果，供上层决定后续流程。
        """
        if suffix in self.SUFFIXES:
            return True
        return mime == "application/pdf"

    def parse(self, data: bytes, *, filename: str = "") -> ParsedDocument:
        """目的：将输入内容转换为统一的内部表示，屏蔽原始格式差异。
        结果：返回标准化解析结果，便于后续链路复用和扩展。
        """
        try:
            from pypdf import PdfReader
        except ImportError as e:
            raise ParseError("需要安装 pypdf：pip install pypdf") from e
        reader = PdfReader(io.BytesIO(data))
        parts: list[str] = []
        for page in reader.pages:
            t = page.extract_text() or ""
            parts.append(t)
        text = "\n".join(parts).strip()
        meta = {"parser": "pdf", "filename": filename, "pages": len(reader.pages)}
        return ParsedDocument(text=text, metadata=meta)


class DocxParser(DocumentParser):
    """目的：封装特定格式的解析逻辑，将原始内容转换为标准文本结果。
    结果：上游流程可以复用统一解析接口处理不同文件类型，降低适配成本。
    """

    # 目的：保存 SUFFIXES 字段，用于 DocxParser 的业务状态、配置或序列化。
    # 结果：实例在读写、校验和协作时可以获得稳定的 SUFFIXES 值。
    SUFFIXES = frozenset({".docx"})

    def supports(self, mime: str | None, suffix: str) -> bool:
        """目的：根据当前输入执行条件判断，统一布尔分支的判定逻辑。
        结果：返回明确的判断结果，供上层决定后续流程。
        """
        if suffix in self.SUFFIXES:
            return True
        return mime in ("application/vnd.openxmlformats-officedocument.wordprocessingml.document",)

    def parse(self, data: bytes, *, filename: str = "") -> ParsedDocument:
        """目的：将输入内容转换为统一的内部表示，屏蔽原始格式差异。
        结果：返回标准化解析结果，便于后续链路复用和扩展。
        """
        try:
            import docx
        except ImportError as e:
            raise ParseError("需要安装 python-docx：pip install python-docx") from e
        document = docx.Document(io.BytesIO(data))
        parts = [p.text for p in document.paragraphs if p.text.strip()]
        text = "\n".join(parts)
        return ParsedDocument(text=text, metadata={"parser": "docx", "filename": filename})


class XlsxParser(DocumentParser):
    """目的：封装特定格式的解析逻辑，将原始内容转换为标准文本结果。
    结果：上游流程可以复用统一解析接口处理不同文件类型，降低适配成本。
    """

    # 目的：保存 SUFFIXES 字段，用于 XlsxParser 的业务状态、配置或序列化。
    # 结果：实例在读写、校验和协作时可以获得稳定的 SUFFIXES 值。
    SUFFIXES = frozenset({".xlsx", ".xlsm"})

    def supports(self, mime: str | None, suffix: str) -> bool:
        """目的：根据当前输入执行条件判断，统一布尔分支的判定逻辑。
        结果：返回明确的判断结果，供上层决定后续流程。
        """
        if suffix in self.SUFFIXES:
            return True
        return mime in ("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",)

    def parse(self, data: bytes, *, filename: str = "") -> ParsedDocument:
        """目的：将输入内容转换为统一的内部表示，屏蔽原始格式差异。
        结果：返回标准化解析结果，便于后续链路复用和扩展。
        """
        try:
            from openpyxl import load_workbook
        except ImportError as e:
            raise ParseError("需要安装 openpyxl：pip install openpyxl") from e
        wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
        lines: list[str] = []
        for sheet in wb:
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    lines.append(" | ".join(cells))
        text = "\n".join(lines)
        return ParsedDocument(text=text, metadata={"parser": "xlsx", "filename": filename})


class PptxParser(DocumentParser):
    """目的：封装特定格式的解析逻辑，将原始内容转换为标准文本结果。
    结果：上游流程可以复用统一解析接口处理不同文件类型，降低适配成本。
    """

    # 目的：保存 SUFFIXES 字段，用于 PptxParser 的业务状态、配置或序列化。
    # 结果：实例在读写、校验和协作时可以获得稳定的 SUFFIXES 值。
    SUFFIXES = frozenset({".pptx"})

    def supports(self, mime: str | None, suffix: str) -> bool:
        """目的：根据当前输入执行条件判断，统一布尔分支的判定逻辑。
        结果：返回明确的判断结果，供上层决定后续流程。
        """
        if suffix in self.SUFFIXES:
            return True
        return mime in (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )

    def parse(self, data: bytes, *, filename: str = "") -> ParsedDocument:
        """目的：将输入内容转换为统一的内部表示，屏蔽原始格式差异。
        结果：返回标准化解析结果，便于后续链路复用和扩展。
        """
        try:
            from pptx import Presentation
        except ImportError as e:
            raise ParseError("需要安装 python-pptx：pip install python-pptx") from e
        prs = Presentation(io.BytesIO(data))
        parts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    parts.append(shape.text.strip())
        text = "\n".join(parts)
        return ParsedDocument(text=text, metadata={"parser": "pptx", "filename": filename})


class RtfParser(DocumentParser):
    """目的：封装特定格式的解析逻辑，将原始内容转换为标准文本结果。
    结果：上游流程可以复用统一解析接口处理不同文件类型，降低适配成本。
    """

    # 目的：保存 SUFFIXES 字段，用于 RtfParser 的业务状态、配置或序列化。
    # 结果：实例在读写、校验和协作时可以获得稳定的 SUFFIXES 值。
    SUFFIXES = frozenset({".rtf"})

    def supports(self, mime: str | None, suffix: str) -> bool:
        """目的：根据当前输入执行条件判断，统一布尔分支的判定逻辑。
        结果：返回明确的判断结果，供上层决定后续流程。
        """
        if suffix in self.SUFFIXES:
            return True
        return mime == "application/rtf"

    def parse(self, data: bytes, *, filename: str = "") -> ParsedDocument:
        """目的：将输入内容转换为统一的内部表示，屏蔽原始格式差异。
        结果：返回标准化解析结果，便于后续链路复用和扩展。
        """
        raw = _decode_best_effort(data)
        # 去掉 {\*...} 与 \word 控制序列，保留可读片段
        text = re.sub(r"\\[a-z]+\d* ?", " ", raw)
        text = re.sub(r"\{[^}]*\}", " ", text)
        text = re.sub(r"\s+", " ", text).strip()
        return ParsedDocument(text=text, metadata={"parser": "rtf_fallback", "filename": filename})


class ZipArchiveParser(DocumentParser):
    """目的：封装特定格式的解析逻辑，将原始内容转换为标准文本结果。
    结果：上游流程可以复用统一解析接口处理不同文件类型，降低适配成本。
    """

    # 目的：保存 SUFFIXES 字段，用于 ZipArchiveParser 的业务状态、配置或序列化。
    # 结果：实例在读写、校验和协作时可以获得稳定的 SUFFIXES 值。
    SUFFIXES = frozenset({".zip"})

    def supports(self, mime: str | None, suffix: str) -> bool:
        """目的：根据当前输入执行条件判断，统一布尔分支的判定逻辑。
        结果：返回明确的判断结果，供上层决定后续流程。
        """
        if suffix in self.SUFFIXES:
            return True
        return mime == "application/zip"

    def parse(self, data: bytes, *, filename: str = "") -> ParsedDocument:
        """目的：将输入内容转换为统一的内部表示，屏蔽原始格式差异。
        结果：返回标准化解析结果，便于后续链路复用和扩展。
        """
        parts: list[str] = []
        try:
            with zipfile.ZipFile(io.BytesIO(data)) as zf:
                for name in sorted(zf.namelist()):
                    if name.endswith("/") or name.startswith("__MACOSX"):
                        continue
                    lower = name.lower()
                    if not any(lower.endswith(ext) for ext in (".md", ".txt", ".markdown")):
                        continue
                    with zf.open(name) as f:
                        parts.append(f"### {name}\n{_decode_best_effort(f.read())}")
        except zipfile.BadZipFile as e:
            raise ParseError("invalid zip") from e
        if not parts:
            raise ParseError("zip 中未找到 .md/.txt")
        text = "\n\n".join(parts)
        meta = {"parser": "zip_text_bundle", "filename": filename}
        return ParsedDocument(text=text, metadata=meta)


class Utf8FallbackParser(DocumentParser):
    """目的：封装特定格式的解析逻辑，将原始内容转换为标准文本结果。
    结果：上游流程可以复用统一解析接口处理不同文件类型，降低适配成本。
    """

    def supports(self, mime: str | None, suffix: str) -> bool:
        """目的：根据当前输入执行条件判断，统一布尔分支的判定逻辑。
        结果：返回明确的判断结果，供上层决定后续流程。
        """
        return True

    def parse(self, data: bytes, *, filename: str = "") -> ParsedDocument:
        """目的：将输入内容转换为统一的内部表示，屏蔽原始格式差异。
        结果：返回标准化解析结果，便于后续链路复用和扩展。
        """
        text = _decode_best_effort(data)
        meta = {
            "parser": "utf8_fallback",
            "filename": filename,
            "warning": "未识别类型，按文本兜底",
        }
        return ParsedDocument(text=text, metadata=meta)


def default_parser_chain() -> list[DocumentParser]:
    """目的：执行返回按优先级排序的策略列表（具体先于兜底）相关逻辑。
    结果：返回当前步骤的处理结果，供后续流程继续使用。
    """
    return [
        PdfParser(),
        DocxParser(),
        XlsxParser(),
        PptxParser(),
        HtmlParser(),
        JsonParser(),
        CsvParser(),
        XmlParser(),
        RtfParser(),
        ZipArchiveParser(),
        PlainTextParser(),
        Utf8FallbackParser(),
    ]
